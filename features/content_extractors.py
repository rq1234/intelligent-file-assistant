"""
Content extraction utilities for different file types
"""
from pathlib import Path

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None


def extract_image_text(file_path):
    """
    Extract text from images using OCR
    
    Args:
        file_path: Path to image file
        
    Returns:
        str: Extracted text or empty string
    """
    if not Image or not pytesseract:
        return ""
    
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text[:500]
    except Exception:
        return ""


def is_image(file_path):
    """Check if file is an image"""
    ext = file_path.lower()
    return ext.endswith(('.png', '.jpg', '.jpeg'))


def extract_title(file_path):
    """
    Extract title from PDF or DOCX files
    
    Args:
        file_path: Path to the file
        
    Returns:
        str: Extracted title or empty string
    """
    if file_path.endswith(".pdf"):
        if PdfReader:
            try:
                reader = PdfReader(file_path)
                if reader.pages:
                    return reader.pages[0].extract_text()[:500]
            except Exception:
                pass
    elif file_path.endswith(".docx"):
        if Document:
            try:
                doc = Document(file_path)
                return " ".join(p.text for p in doc.paragraphs[:5])
            except Exception:
                pass
    elif file_path.endswith(".pptx"):
        if Presentation:
            try:
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides[:5]:
                    for shape in slide.shapes:
                        text = getattr(shape, "text", "")
                        if text:
                            texts.append(text)
                return " ".join(texts)[:500]
            except Exception:
                pass
    elif file_path.endswith(".ppt"):
        # Old PowerPoint format - content extraction not supported
        # Return a message so LLM knows this is a PowerPoint file
        return "[PowerPoint file - content not extractable, old .ppt format]"
    return ""


def extract_text_content(file_path, max_length=1000):
    """
    Extract text content from a file
    
    Args:
        file_path: Path to the file
        max_length: Maximum length of text to extract
        
    Returns:
        str: Extracted text content
    """
    path = Path(file_path)
    extension = path.suffix.lower().lstrip('.')
    
    try:
        if extension in ['txt', 'md', 'log', 'csv']:
            return extract_plain_text(path, max_length)
        elif extension in ['pdf']:
            return extract_pdf_text(path, max_length)
        elif extension in ['doc', 'docx']:
            return extract_doc_text(path, max_length)
        else:
            return ""
    except Exception as e:
        print(f"Error extracting content from {path}: {e}")
        return ""


def extract_plain_text(file_path, max_length=1000):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read(max_length)
        return content
    except Exception as e:
        print(f"Error reading plain text: {e}")
        return ""


def extract_pdf_text(file_path, max_length=1000):
    """Extract text from PDF files"""
    if not PdfReader:
        return ""
    
    try:
        reader = PdfReader(file_path)
        if reader.pages:
            text = reader.pages[0].extract_text()
            return text[:max_length]
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
    return ""


def extract_doc_text(file_path, max_length=1000):
    """Extract text from Word documents"""
    if not Document:
        return ""
    
    try:
        doc = Document(file_path)
        text = " ".join(p.text for p in doc.paragraphs)
        return text[:max_length]
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
    return ""
